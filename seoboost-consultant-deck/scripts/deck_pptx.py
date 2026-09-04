#!/usr/bin/env python3
"""
deck_pptx.py — render a consultant-grade business-case deck (PPTX) in the house style.

Reads a deck spec (JSON, see ../assets/deck.example.json) and emits a .pptx built from a
small set of slide archetypes: cover, divider, content (action-title + bullets), table
(options comparison with RAG cells), conclusion. One design system, editable in PowerPoint.

Usage:
    python3 deck_pptx.py --spec ../assets/deck.example.json --out /tmp/deck.pptx
    python3 deck_pptx.py --spec my_deck.json --out out.pptx

Requires: python-pptx  (pip install python-pptx)

Design tokens live in THEME below — change once to re-skin. Keep in sync with
../references/design-system.md and deck_html.py.
"""
import argparse
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

# --- Design tokens (defaults; override per-deck via spec['theme']) --------------------
THEME = {
    "navy": "1B3A6B", "navy_deep": "12285A", "red": "C00000", "teal": "00B0A0",
    "ink": "262626", "grey": "7F7F7F", "grey_band": "8496B0", "paper": "FFFFFF",
    "wash": "EEF2F8", "green": "2E7D32", "amber": "ED9B00",
}
FONT = "Segoe UI"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#").upper())


def theme(spec):
    t = dict(THEME)
    for k, v in (spec.get("theme") or {}).items():
        t[k] = v.lstrip("#").upper()
    return t


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def _text(tf, runs, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
          font=FONT, space_after=6):
    """runs: str or list[str] (each str = its own paragraph/bullet)."""
    if isinstance(runs, str):
        runs = [runs]
    tf.word_wrap = True
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        f = r.font
        f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
        f.color.rgb = color
    return tf


def _rect(slide, l, t, w, h, fill_hex, line_hex=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        sp.line.color.rgb = rgb(line_hex)
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --- Archetypes -----------------------------------------------------------------------
def furniture(slide, t, page):
    """Top rule + red marker + page number — every content/table slide."""
    _rect(slide, 0.55, 0.42, 0.22, 0.22, t["red"])              # red marker square
    _rect(slide, 0.85, 0.5, 11.9, 0.03, t["navy"])              # navy top rule
    if page is not None:
        b = _box(slide, 0.55, 7.05, 1.0, 0.35)
        _text(b.text_frame, str(page), 9, rgb(t["grey"]))


def footnotes(slide, t, notes):
    if not notes:
        return
    b = _box(slide, 0.55, 6.55, 12.2, 0.5)
    _text(b.text_frame, notes, 8, rgb(t["grey"]), space_after=1)


def slide_cover(prs, t, s):
    sl = blank(prs)
    _rect(sl, 0, 0, 13.333, 4.1, t["wash"])                    # hero placeholder band
    hb = _box(sl, 0.6, 1.6, 12, 1)
    _text(hb.text_frame, "[ hero image ]", 12, rgb(t["grey"]), align=PP_ALIGN.CENTER)
    _rect(sl, 0, 4.1, 13.333, 0.06, t["red"])                  # red rule
    _rect(sl, 0, 4.16, 13.333, 3.34, t["navy"])                # navy band
    tb = _box(sl, 0.7, 4.5, 9, 1.6)
    _text(tb.text_frame, s.get("title", ""), 32, rgb(t["paper"]), bold=True)
    if s.get("subtitle"):
        _text(tb.text_frame.add_paragraph() and tb.text_frame, "", 1, rgb(t["paper"]))  # noop spacer
    sb = _box(sl, 0.7, 5.4, 9, 0.8)
    _text(sb.text_frame, s.get("subtitle", ""), 18, rgb(t["paper"]), bold=True)
    fb = _box(sl, 0.7, 6.5, 9, 0.8)
    _text(fb.text_frame, [s.get("tag", "Final Presentation"), s.get("date", "")],
          13, rgb(t["paper"]), space_after=2)
    return sl


def slide_divider(prs, t, s):
    sl = blank(prs)
    _rect(sl, 0, 0, 13.333, 7.5, t["navy_deep"])
    _rect(sl, 0.7, 3.05, 0.9, 0.12, t["red"])
    tb = _box(sl, 0.7, 3.2, 11, 1.4)
    label = f'{s.get("letter","")}. {s.get("title","")}'.strip(". ")
    _text(tb.text_frame, label, 38, rgb(t["paper"]), bold=True)
    return sl


def slide_content(prs, t, s):
    sl = blank(prs)
    furniture(sl, t, s.get("page"))
    if s.get("kicker"):
        kb = _box(sl, 0.85, 0.75, 11, 0.4)
        _text(kb.text_frame, s["kicker"], 13, rgb(t["red"]), bold=True)
    tb = _box(sl, 0.85, 1.15, 11.7, 1.2)
    _text(tb.text_frame, s.get("title", ""), 26, rgb(t["navy"]), bold=True)
    bb = _box(sl, 0.95, 2.5, 11.5, 3.7)
    bullets = s.get("bullets", [])
    _text(bb.text_frame, [f"•  {x}" for x in bullets], 16, rgb(t["ink"]), space_after=10)
    footnotes(sl, t, s.get("footnotes"))
    return sl


def slide_table(prs, t, s):
    sl = blank(prs)
    furniture(sl, t, s.get("page"))
    if s.get("kicker"):
        kb = _box(sl, 0.85, 0.75, 11, 0.4)
        _text(kb.text_frame, s["kicker"], 13, rgb(t["red"]), bold=True)
    tb = _box(sl, 0.85, 1.15, 11.7, 1.1)
    _text(tb.text_frame, s.get("title", ""), 24, rgb(t["navy"]), bold=True)

    cols, rows = s["columns"], s["rows"]
    hl = s.get("highlight_col")
    rag = {"green": t["green"], "amber": t["amber"], "red": t["red"]}
    n_r, n_c = len(rows) + 1, len(cols)
    gt = sl.shapes.add_table(n_r, n_c, Inches(0.85), Inches(2.5),
                             Inches(11.6), Inches(0.5 * n_r)).table
    for c, name in enumerate(cols):
        cell = gt.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(t["red"] if (hl and c == hl) else t["navy"])
        _text(cell.text_frame, name, 12, rgb(t["paper"]), bold=True)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            key = str(val).lower()
            if key in rag:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(rag[key])
                _text(cell.text_frame, key.upper(), 11, rgb(t["paper"]), bold=True,
                      align=PP_ALIGN.CENTER)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(t["wash"] if r % 2 else t["paper"])
                _text(cell.text_frame, str(val), 12, rgb(t["ink"]))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    footnotes(sl, t, s.get("footnotes"))
    return sl


def slide_conclusion(prs, t, s):
    sl = blank(prs)
    furniture(sl, t, s.get("page"))
    kb = _box(sl, 0.85, 0.75, 11, 0.4)
    _text(kb.text_frame, "Conclusion", 13, rgb(t["red"]), bold=True)
    tb = _box(sl, 0.85, 1.15, 11.7, 1.2)
    _text(tb.text_frame, s.get("title", ""), 26, rgb(t["navy"]), bold=True)
    bb = _box(sl, 0.95, 2.5, 11.5, 2.2)
    _text(bb.text_frame, [f"•  {x}" for x in s.get("bullets", [])], 15, rgb(t["ink"]),
          space_after=8)
    if s.get("next_steps"):
        nb = _box(sl, 0.95, 4.9, 11.5, 1.5)
        _text(nb.text_frame, "Next steps", 13, rgb(t["red"]), bold=True)
        _text(nb.text_frame.add_paragraph() and nb.text_frame, "", 1, rgb(t["paper"]))
        for x in s["next_steps"]:
            p = nb.text_frame.add_paragraph()
            r = p.add_run(); r.text = f"→  {x}"
            r.font.size, r.font.name = Pt(14), FONT
            r.font.color.rgb = rgb(t["ink"])
    footnotes(sl, t, s.get("footnotes"))
    return sl


RENDERERS = {
    "cover": slide_cover, "divider": slide_divider, "content": slide_content,
    "table": slide_table, "conclusion": slide_conclusion,
}


def build(spec, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    t = theme(spec)
    for s in spec.get("slides", []):
        fn = RENDERERS.get(s.get("type"))
        if not fn:
            print(f"  ! unknown slide type: {s.get('type')!r} — skipped", file=sys.stderr)
            continue
        fn(prs, t, s)
    prs.save(out)
    print(f"Wrote {out} — {len(prs.slides._sldIdLst)} slides")


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--spec", required=True, help="deck spec JSON (see assets/deck.example.json)")
    ap.add_argument("--out", default="deck.pptx", help="output .pptx path")
    a = ap.parse_args()
    with open(a.spec, encoding="utf-8") as f:
        spec = json.load(f)
    build(spec, a.out)


if __name__ == "__main__":
    main()
